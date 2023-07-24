import PyPDF2
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor



def extract_all_info_opq(opq_file):
	pdf_reader = PyPDF2.PdfReader(opq_file)
	text = ""
	for page in range(len(pdf_reader.pages)):
		page_obj = pdf_reader.pages[page]
		text += page_obj.extract_text()

	marker = "1. Leda och Ta beslut 12345"

	lines = text.split("\n")
	index = next((i for i, line in enumerate(lines) if marker in line), None)

	if index is not None:
		trimmed_string = "\n".join(lines[index:])
	
	index = trimmed_string.find(
		'1.1 Besluta  och starta aktiviteterTar ansvar för åtgärder, projekt och  personal; tar initiativ  och arbetar')

	if index != -1:
		return trimmed_string[:index]
	else:
		return trimmed_string

def extract_all_verify(verify_file):
	pdf_reader = PyPDF2.PdfReader(verify_file)
	text = ""
	for page in range(len(pdf_reader.pages)):
		page_obj = pdf_reader.pages[page]
		text += page_obj.extract_text()

	return text

def string_to_list_opq(opq_file):
	list_of_words = ['DeDN', 'Spacer', '1', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '2.2', '3.1', '3.2', '3.3', '4.1',
	                 '4.2', '4.3', '5.1', '5.2', '5.3', '6.1', '6.2', '6.3', '7.1', '7.2', '8.1', '8.2', '\uf0fc\uf0fc',
	                 '\uf0fc', '\uf06c\uf06c', '\uf06c', '\uf0fb\uf0fb', '\uf0fb', '©']
	long_string = extract_all_info_opq(opq_file)
	len_list = long_string.split("\n")
	list = long_string.split("\n")
	remove_list = []
	final_list = []

	for i in range(len(len_list)):

		if not any(list[i].startswith(word) for word in list_of_words):
			joined_string = list[i - 1] + list[i]
			list[i - 1] = joined_string
			remove_list.append(list[i])

	for word in list:
		if word not in remove_list:
			final_list.append(word)

	modified_list = []

	for string in final_list:
		while "  " in string:
			string = string.replace("  ", " ")
		modified_list.append(string)

	return modified_list

def string_to_list_verify(verify_file):
	long_string = extract_all_verify(verify_file)
	long_list = long_string.split("\n")

	name = long_list[1]
	generell = long_list[12].split('G+')[1]
	induktiv = long_list[13].split('Induktivt resonemang')[1]
	numerisk = long_list[14].split('Numeriska färdigheter')[1]
	deduktiv = long_list[15].split('Deduktivt resonemang')[1]

	induktiv_text_list = long_list[38:41]
	numerisk_text_list = long_list[53:56]
	deduktiv_text_list = long_list[64:70]

	induktiv_string = ''.join(induktiv_text_list)
	numerisk_string = ''.join(numerisk_text_list)
	deduktiv_string = ''.join(deduktiv_text_list)

	all_percent_list = [generell, induktiv, numerisk, deduktiv]
	all_percent_list = [s.replace(' ', '') for s in all_percent_list]
	all_text_list = [induktiv_string, numerisk_string, deduktiv_string]

	modified_list = []

	for string in all_text_list:
		while "  " in string:
			string = string.replace("  ", " ")
		modified_list.append(string)

	return all_percent_list, modified_list, name

def get_dimmension(index, number, opq_file):
	chosen_list = []
	full_list = string_to_list_opq(opq_file)
	dimmension = full_list[index].split(" ", 1)
	chosen_list.append(dimmension[1])

	for n in range(index + 1, index + number + 1):
		a = full_list[n].split(" ", 1)
		chosen_list.append(a)

	return chosen_list

def get_chosen_dimmension_list(opq_file, checkbox_list):
	list_of_dimmenssions = [['Besluta och starta aktiviteter', 1, 4], ['Leda och följa upp', 6, 4],
	                        ['Arbeta med människor', 12, 5], ['Stå fast vid principer och värderingar', 18, 2],
	                        ['Skapa relationer och nätverk', 22, 4], ['Övertala och påverka', 27, 5],
	                        ['Presentera och kommunicera information', 33, 4], ['Skriva och rapportera', 41, 4],
	                        ['Tillämpa expertis och teknologi', 46, 3], ['Analysera', 50, 3], ['Lära och utforska', 55, 4],
	                        ['Skapa och uppfinna', 60, 4], ['Formulera strategier och koncept', 65, 4],
	                        ['Planera och organisera', 72, 4], ['Leverera resultat och uppfylla kundförväntningar', 77, 4],
	                        ['Följa instruktioner och procedurer', 82, 3], ['Anpassa och reagera på förändring', 87, 4],
	                        ['Hantera krav och motgångar', 92, 4], ['Uppnå personliga arbetsmål', 99, 4],
	                        ['Företagaranda och kommersiellt  tänkande', 104, 3]]

	string_list = []
	index_list = []
	for i in range(len(list_of_dimmenssions)):
		d = get_dimmension(list_of_dimmenssions[i][1], list_of_dimmenssions[i][2], opq_file)
		if list_of_dimmenssions[i][0] in checkbox_list:
			index_list.append(list_of_dimmenssions[i][2])
			string_list.append(d)

	return string_list, index_list

def make_slide_8_to_12(presentation, opq_file, checkbox_list):
	slide_index = 7
	dimmension_list, index_list = get_chosen_dimmension_list(opq_file, checkbox_list)
	only_dimension_list = checkbox_list

	for i in range(len(only_dimension_list)):

		# Ensure the slide index is valid
		if slide_index < len(presentation.slides):
			slide = presentation.slides[slide_index]

			textbox_index = 0  # Modify the first textbox
			if len(slide.shapes) > textbox_index:
				shape = slide.shapes[textbox_index]
				if shape.has_text_frame:
					text_frame = shape.text_frame

					# Modify the text content
					text_frame.text = only_dimension_list[i]

					# Modify the font size
					for paragraph in text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(24)
							run.font.bold = False

			textbox_index = 1  # Modify the second textbox
			if len(slide.shapes) > textbox_index:
				shape = slide.shapes[textbox_index]
				if shape.has_text_frame:
					text_frame = shape.text_frame

					# Modify the text content
					text_frame.text = only_dimension_list[i]

					# Modify the font size
					for paragraph in text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(16)
							run.font.bold = False

							# Modify the font color
							font_color = RGBColor(0, 0, 0)  # Customize the RGB color values as needed
							run.font.color.rgb = font_color

			# Iterate through the shapes to find the table
			for shape in slide.shapes:
				if shape.has_table:
					table = shape.table

					# Check if the table has at least four rows and two columns
					if len(table.rows) >= 5 and len(table.columns) > 1:
						# Modify the first and second columns
						texts_first_column = []
						texts_second_column = []
						for b in range(index_list[i]):
							try:
								texts_first_column.append(dimmension_list[i][b+1][0])
								texts_second_column.append(dimmension_list[i][b+1][1])
							except Exception:
								continue

						for n in range(index_list[i]):  # Modify the current row and the next three rows
							# Modify cell in the first column
							cell_first_column = table.cell(n, 0)
							if cell_first_column.text:
								text_frame_first_column = cell_first_column.text_frame
								paragraph_first_column = text_frame_first_column.paragraphs[0]
								run_first_column = paragraph_first_column.runs[0]

								run_first_column.text = texts_first_column[n]
								run_first_column.font.bold = False
								run_first_column.font.size = Pt(14)

								if text_frame_first_column.text.startswith('\uf0fc'):
									run_first_column.font.color.rgb = RGBColor(0, 255, 0)
								elif text_frame_first_column.text.startswith('\uf0fb'):
									run_first_column.font.color.rgb = RGBColor(255, 0, 0)
								else:
									run_first_column.font.color.rgb = RGBColor(0, 0, 0)

							# Modify cell in the second column
							cell_second_column = table.cell(n, 1)
							if cell_second_column.text:
								text_frame_second_column = cell_second_column.text_frame
								paragraph_second_column = text_frame_second_column.paragraphs[0]
								run_second_column = paragraph_second_column.runs[0]

								run_second_column.text = texts_second_column[n]
								run_second_column.font.bold = False
								run_second_column.font.size = Pt(14)
								run_second_column.font.color.rgb = RGBColor(0, 0, 0)
		slide_index = 8 + i
	
	presentation.save()

def make_slide_4_and_7_OPQ(presentation, slide_index, font_size, checkbox_list):

	if slide_index < len(presentation.slides):
		slide = presentation.slides[slide_index]

		# Method: Find the table based on cell values
		target_table_cells = [
			["Test"],
			["Test"],
			["Test"],
			["Test"],
			["Test"],
		]
		new_texts = checkbox_list

		for shape in slide.shapes:
			if shape.has_table:
				table = shape.table

				# Check if the table has the desired number of rows and columns
				if len(table.rows) >= len(target_table_cells) and len(table.columns) >= 1:
					found_table = True

					# Check if the cell values match the target_table_cells
					for row_index, row_cells in enumerate(target_table_cells):
						for col_index, cell_value in enumerate(row_cells):
							cell = table.cell(row_index, col_index)
							if cell.text != cell_value:
								found_table = False
								break

					if found_table:
						# Update the table with new texts
						for i in range(len(new_texts)):
							cell = table.cell(i, 0)
							if cell.text:
								text_frame = cell.text_frame
								paragraph = text_frame.paragraphs[0]
								run = paragraph.runs[0]

								run.text = new_texts[i]
								run.font.size = Pt(font_size)
								run.font.bold = False

	# Save the modified presentation
	presentation.save()

def make_slide_4_and_7_verify(presentation, slide_index, font_size, verify_file):
    

    if slide_index < len(presentation.slides):
        slide = presentation.slides[slide_index]
        target_table_cells = [
            ["Test"],
            ["Test"],
            ["Test"],
            ["Test"],
            ["Test"],
        ]
        percent_list, text_list, name = string_to_list_verify(verify_file)

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table

                if len(table.rows) >= len(target_table_cells) and len(table.columns) > 1:
                    found_table = True

                    # Check if the cell values match the target_table_cells
                    for row_index, row_cells in enumerate(target_table_cells):
                        for col_index, cell_value in enumerate(row_cells):
                            if (
                                row_index < len(table.rows) and
                                col_index < len(table.columns)
                            ):
                                cell = table.cell(row_index, col_index)
                                if cell.text != cell_value:
                                    found_table = False
                                    break
                            else:
                                found_table = False
                                break

                    if found_table:
                        # Update the table with new texts
                        for i in range(len(text_list)):
                            if i < len(table.rows) and i < len(text_list):
                                cell = table.cell(i, 0)
                                if cell.text:
                                    text_frame = cell.text_frame
                                    paragraph = text_frame.paragraphs[0]
                                    run = paragraph.runs[0]

                                    run.text = text_list[i]
                                    run.font.size = Pt(font_size)
                                    run.font.bold = False

                        # Update the table with the name
                        if len(table.columns) > 1 and len(table.rows) > 5:
                            cell = table.cell(5, 1)
                            if cell.text:
                                text_frame = cell.text_frame
                                paragraph = text_frame.paragraphs[0]
                                run = paragraph.runs[0]

                                run.text = name
                                run.font.size = Pt(font_size)
                                run.font.bold = False

    # Save the modified presentation
    presentation.save()



def run(opq_file, verify_file, checkbox_list):
    modified_presentation_path = "modified_presentation.pptx"
    presentation = Presentation("Fördjupad bedömning - mall.pptx")

    # Make modifications to the presentation
    make_slide_8_to_12(presentation, opq_file, checkbox_list)
    make_slide_4_and_7_OPQ(presentation, 3, 24, checkbox_list)
    make_slide_4_and_7_verify(presentation, 6, 24, verify_file)

    # Save the modified presentation
    presentation.save(modified_presentation_path)

    return modified_presentation_path


	
    
